"""
Ce fichier normalize les données dans <cellvision_root>/datasets/ :
- récupération des images des .pptx
- renommage des .h5
- renommage des .jpg

"""

import os
import shutil
import re
from pathlib import Path
from pptx import Presentation

def extract_images_from_pptx(pptx_path, output_folder):
    """Extrait toutes les images d'un fichier .pptx ou .ppt vers le dossier de sortie."""
    try:
        prs = Presentation(pptx_path)
        img_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "image"):
                    continue

                image = shape.image
                
                if image.ext.lower() in ['jpg', 'jpeg']:
                    img_count += 1
                    image_filename = f"extracted {img_count}.jpg"
                    with open(output_folder / image_filename, 'wb') as f:
                        f.write(image.blob)
    except Exception as e:
        print(f"Erreur lors de l'extraction de {pptx_path}: {e}")

def normalized_filename(folder_path, filename):
    return (
        filename
            .replace(folder_path.name, "") # Retirer la duplication de date
            .replace("_","")
            .replace("-","")
            .strip(" ")
    )

def normalize_datasets(input_root, output_root):
    input_path = Path(input_root)
    output_path = Path(output_root)
    
    shutil.rmtree(output_path, ignore_errors=True)
    output_path.mkdir(parents=True, exist_ok=True)

    # Regex pour trouver les dossiers nommés par 8 chiffres (YYYYMMDD)
    date_pattern = re.compile(r'^(\d{4})(\d{2})(\d{2})$')

    for source_folder in input_path.iterdir():
        if not source_folder.is_dir():
            continue

        match = date_pattern.match(source_folder.name)
        
        if not match:
            continue

        # Formater la date : 20210303 -> 2021-03-03
        year, month, day = match.groups()
        formatted_date = f"{year}-{month}-{day}"
        
        target_folder = output_path / formatted_date
        target_folder.mkdir(exist_ok=True)

        print(f"Traitement du dossier : {source_folder.name} -> {formatted_date}")

        all_recursive_files = source_folder.rglob('*')
        
        for file_path in all_recursive_files:
            if not file_path.is_file():
                continue
            
            if("mask" in file_path.name.lower()):
                continue

            ext = file_path.suffix.lower()
            
            if ext in ['.jpg', '.jpeg', '.h5']:
                shutil.copy2(file_path, target_folder / normalized_filename(source_folder, file_path.name))

            elif ext in ['.pptx']:
                extract_images_from_pptx(file_path, target_folder)

if __name__ == "__main__":
    # Configuration des dossiers
    SOURCE = "datasets"
    DESTINATION = "normalized_datasets"
    
    normalize_datasets(SOURCE, DESTINATION)
    print("\nNormalisation terminée !")